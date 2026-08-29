from __future__ import annotations

import hashlib
import io
import json
import logging
import re
import uuid
import zipfile
from dataclasses import dataclass
from pathlib import Path, PurePosixPath
from defusedxml import ElementTree as DefusedET
from xml.sax.saxutils import escape  # nosec B406: serializer de XML gerado, não parser de entrada

from sqlalchemy.orm import Session

from ..config import Settings
from ..models import Artifact
from .blob_storage import BlobStorage, BlobStorageError, build_blob_storage

DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
TXT_MIME = "text/plain"
PDF_MIME = "application/pdf"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
MARKDOWN_MIME = "text/markdown"
JSON_MIME = "application/json"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

artifact_runtime_logger = logging.getLogger("uvicorn.error")

_ARTIFACT_INTENT_RE = re.compile(
    r"\b(?:gere|gerar|crie|criar|exporte|exportar|salve|salvar|produza|produzir|create|generate|export|save)\b",
    re.I,
)
_DOCX_RE = re.compile(r"(?:\.docx?\b|\bdocx?\b|\bword\b)", re.I)
_TXT_RE = re.compile(r"(?:\.txt\b|\btxt\b|\btexto simples\b|\bplain text\b)", re.I)
_PDF_RE = re.compile(r"(?:\.pdf\b|\bpdf\b)", re.I)
_PPTX_RE = re.compile(r"(?:\.pptx?\b|\bpptx?\b|\bpowerpoint\b|\bapresenta(?:ção|cao)\b)", re.I)
_MARKDOWN_RE = re.compile(r"(?:\.md\b|\bmarkdown\b)", re.I)
_JSON_RE = re.compile(r"(?:\.json\b|\bjson\b)", re.I)
_XLSX_RE = re.compile(r"(?:\.xlsx?\b|\bxlsx?\b|\bexcel\b|\bplanilha\b)", re.I)


class ArtifactGenerationError(RuntimeError):
    code = "ARTIFACT_GENERATION_ERROR"


class ArtifactFormatUnsupported(ArtifactGenerationError):
    code = "ARTIFACT_FORMAT_UNSUPPORTED"


class ArtifactValidationFailed(ArtifactGenerationError):
    code = "ARTIFACT_VALIDATION_FAILED"


class ArtifactStorageError(ArtifactGenerationError):
    code = "ARTIFACT_STORAGE_ERROR"


@dataclass(frozen=True, slots=True)
class ArtifactIntent:
    requested_format: str
    extension: str
    mime_type: str


@dataclass(frozen=True, slots=True)
class ValidatedArtifactBytes:
    filename: str
    mime_type: str
    data: bytes
    sha256: str
    semantic_text: str
    renderer: str


@dataclass(frozen=True, slots=True)
class PersistedArtifact:
    artifact: Artifact
    download_path: str
    provenance_path: str


def detect_artifact_intent(message: str) -> ArtifactIntent | None:
    text = (message or "").strip()
    if not text or not _ARTIFACT_INTENT_RE.search(text):
        return None
    if _DOCX_RE.search(text):
        return ArtifactIntent("docx", ".docx", DOCX_MIME)
    if _PDF_RE.search(text):
        return ArtifactIntent("pdf", ".pdf", PDF_MIME)
    if _PPTX_RE.search(text):
        return ArtifactIntent("pptx", ".pptx", PPTX_MIME)
    if _MARKDOWN_RE.search(text):
        return ArtifactIntent("markdown", ".md", MARKDOWN_MIME)
    if _JSON_RE.search(text):
        return ArtifactIntent("json", ".json", JSON_MIME)
    if _XLSX_RE.search(text):
        return ArtifactIntent("xlsx", ".xlsx", XLSX_MIME)
    if _TXT_RE.search(text):
        return ArtifactIntent("txt", ".txt", TXT_MIME)
    return None


def artifact_generation_system_message(intent: ArtifactIntent) -> dict[str, str]:
    return {
        "role": "system",
        "content": (
            "ARTIFACT CAPABILITY AVAILABLE FOR THIS TURN. "
            f"The runtime will render and persist the final answer as {intent.requested_format.upper()} "
            "after validation. Produce only the complete document body that should be placed in the file. "
            "Do not narrate execution status. Do not claim that the file is generated, persisted, ready, "
            "or available for download; only the runtime may confirm those states after persistence. "
            "Do not claim that file generation is unavailable. Do not invent a download URL."
        ),
    }


def _safe_filename(name: str, extension: str) -> str:
    base = PurePosixPath((name or "").replace("\\", "/")).name.strip()
    if not base or base in {".", ".."}:
        base = "documento"
    base = re.sub(r"[^A-Za-z0-9._ -]+", "_", base).strip(" ._") or "documento"
    stem = base.rsplit(".", 1)[0] if "." in base else base
    return f"{stem[:120]}{extension}"


def default_filename(intent: ArtifactIntent, *, agent_name: str) -> str:
    label = re.sub(r"[^A-Za-z0-9_-]+", "-", (agent_name or "patroai").strip()).strip("-").lower()
    return _safe_filename(f"patroai-{label or 'artifact'}", intent.extension)


def _docx_bytes(text: str) -> bytes:
    paragraphs = [line.rstrip() for line in (text or "").replace("\r\n", "\n").split("\n")]
    if not any(p.strip() for p in paragraphs):
        raise ArtifactValidationFailed("ARTIFACT_EMPTY_CONTENT")

    body = []
    for paragraph in paragraphs:
        if not paragraph:
            body.append("<w:p/>")
            continue
        body.append(
            "<w:p><w:r><w:t xml:space=\"preserve\">"
            + escape(paragraph)
            + "</w:t></w:r></w:p>"
        )
    body.append(
        '<w:sectPr><w:pgSz w:w="12240" w:h="15840"/>'
        '<w:pgMar w:top="1440" w:right="1440" w:bottom="1440" w:left="1440" '
        'w:header="720" w:footer="720" w:gutter="0"/></w:sectPr>'
    )
    document_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">'
        "<w:body>" + "".join(body) + "</w:body></w:document>"
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/word/document.xml" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>'
        "</Types>"
    )
    rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
        'Target="word/document.xml"/>'
        "</Relationships>"
    )

    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as z:
        z.writestr("[Content_Types].xml", content_types)
        z.writestr("_rels/.rels", rels)
        z.writestr("word/document.xml", document_xml)
    return out.getvalue()


def _strip_markdown_fence(content: str) -> str:
    lines = (content or "").strip().splitlines()
    if lines and lines[0].strip().startswith("```"):
        lines = lines[1:]
    if lines and lines[-1].strip() == "```":
        lines = lines[:-1]
    return "\n".join(lines).strip()


def _json_bytes(content: str) -> bytes:
    source = _strip_markdown_fence(content)
    try:
        parsed = json.loads(source)
    except json.JSONDecodeError as exc:
        raise ArtifactValidationFailed("ARTIFACT_JSON_INVALID") from exc
    return json.dumps(parsed, ensure_ascii=False, indent=2).encode("utf-8")


def _pdf_bytes(text: str) -> bytes:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.units import mm
        from reportlab.pdfgen.canvas import Canvas
    except ImportError as exc:
        raise ArtifactFormatUnsupported("ARTIFACT_PDF_RENDERER_UNAVAILABLE") from exc

    output = io.BytesIO()
    canvas = Canvas(output, pagesize=A4)
    width, height = A4
    left = 20 * mm
    top = height - 20 * mm
    bottom = 20 * mm
    leading = 14
    y = top
    for line in (text or "").replace("\r\n", "\n").split("\n"):
        if y < bottom:
            canvas.showPage()
            y = top
        canvas.drawString(left, y, line[:1800])
        y -= leading
    canvas.save()
    return output.getvalue()


def _pptx_bytes(text: str) -> bytes:
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError as exc:
        raise ArtifactFormatUnsupported("ARTIFACT_PPTX_RENDERER_UNAVAILABLE") from exc

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(12.1), Inches(6.2))
    frame = box.text_frame
    frame.word_wrap = True
    frame.text = (text or "").strip()
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _xlsx_separator_row(line: str) -> bool:
    clean = (line or "").strip()
    if "|" not in clean:
        return False
    cells = [cell.strip() for cell in clean.strip("|").split("|")]
    return bool(cells) and all(
        set(cell.replace("-", "").replace(":", "").strip()) == set()
        for cell in cells
    )


def _xlsx_bytes(text: str) -> bytes:
    try:
        from openpyxl import Workbook
    except ImportError as exc:
        raise ArtifactFormatUnsupported("ARTIFACT_XLSX_RENDERER_UNAVAILABLE") from exc

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "PatroAI"
    lines = (text or "").replace("\r\n", "\n").split("\n")
    rows: list[list[str]] = []
    for line in lines:
        clean = line.strip()
        if not clean:
            continue
        if "|" in clean:
            cells = [cell.strip() for cell in clean.strip("|").split("|")]
            if _xlsx_separator_row(clean):
                continue
            rows.append(cells)
        elif "\t" in clean:
            rows.append([cell.strip() for cell in clean.split("\t")])
        else:
            rows.append([clean])
    if not rows:
        raise ArtifactValidationFailed("ARTIFACT_EMPTY_CONTENT")
    for row in rows:
        sheet.append(row)
    for column in sheet.columns:
        width = min(max(len(str(cell.value or "")) for cell in column) + 2, 60)
        sheet.column_dimensions[column[0].column_letter].width = width
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def _extract_xlsx_text(data: bytes) -> str:
    try:
        from openpyxl import load_workbook
        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        values: list[str] = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(value) for value in row if value is not None and str(value).strip()]
                if cells:
                    values.append(" | ".join(cells))
        workbook.close()
    except Exception as exc:
        raise ArtifactValidationFailed("ARTIFACT_XLSX_VALIDATION_FAILED") from exc
    text = "\n".join(values).strip()
    if not text:
        raise ArtifactValidationFailed("ARTIFACT_XLSX_SEMANTIC_EMPTY")
    return text


def _extract_pdf_text(data: bytes) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data), strict=True)
        return "\n".join(page.extract_text() or "" for page in reader.pages).strip()
    except Exception as exc:
        raise ArtifactValidationFailed("ARTIFACT_PDF_VALIDATION_FAILED") from exc


def _extract_pptx_text(data: bytes) -> str:
    try:
        from pptx import Presentation
        presentation = Presentation(io.BytesIO(data))
        values: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    values.append(shape.text or "")
        text = "\n".join(values).strip()
    except Exception as exc:
        raise ArtifactValidationFailed("ARTIFACT_PPTX_VALIDATION_FAILED") from exc
    if not text:
        raise ArtifactValidationFailed("ARTIFACT_PPTX_SEMANTIC_EMPTY")
    return text


def _extract_json_text(data: bytes) -> str:
    try:
        parsed = json.loads(data.decode("utf-8"))
    except (UnicodeDecodeError, json.JSONDecodeError) as exc:
        raise ArtifactValidationFailed("ARTIFACT_JSON_VALIDATION_FAILED") from exc
    return json.dumps(parsed, ensure_ascii=False, indent=2)


def _extract_docx_text(data: bytes) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as z:
            if z.testzip() is not None:
                raise ArtifactValidationFailed("ARTIFACT_DOCX_CRC_FAILED")
            required = {"[Content_Types].xml", "_rels/.rels", "word/document.xml"}
            if not required.issubset(set(z.namelist())):
                raise ArtifactValidationFailed("ARTIFACT_DOCX_STRUCTURE_INVALID")
            raw = z.read("word/document.xml")
    except ArtifactValidationFailed:
        raise
    except Exception as exc:
        raise ArtifactValidationFailed("ARTIFACT_DOCX_OPEN_FAILED") from exc

    try:
        root = DefusedET.fromstring(
            raw,
            forbid_dtd=True,
            forbid_entities=True,
            forbid_external=True,
        )
    except DefusedET.ParseError as exc:
        raise ArtifactValidationFailed("ARTIFACT_DOCX_XML_INVALID") from exc
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    text = "\n".join((node.text or "") for node in root.iter(f"{ns}t")).strip()
    if not text:
        raise ArtifactValidationFailed("ARTIFACT_DOCX_SEMANTIC_EMPTY")
    return text


def render_and_validate(
    *,
    intent: ArtifactIntent,
    content: str,
    filename: str,
) -> ValidatedArtifactBytes:
    filename = _safe_filename(filename, intent.extension)
    normalized = (content or "").strip()
    if not normalized:
        raise ArtifactValidationFailed("ARTIFACT_EMPTY_CONTENT")

    if intent.requested_format == "docx":
        data = _docx_bytes(normalized)
        semantic_text = _extract_docx_text(data)
        renderer = "patroai_docx_minimal_v1"
    elif intent.requested_format == "pdf":
        data = _pdf_bytes(normalized)
        semantic_text = _extract_pdf_text(data)
        renderer = "patroai_pdf_reportlab_v1"
    elif intent.requested_format == "pptx":
        data = _pptx_bytes(normalized)
        semantic_text = _extract_pptx_text(data)
        renderer = "patroai_pptx_python_v1"
    elif intent.requested_format == "markdown":
        data = normalized.encode("utf-8")
        semantic_text = data.decode("utf-8").strip()
        renderer = "patroai_markdown_v1"
    elif intent.requested_format == "json":
        data = _json_bytes(normalized)
        semantic_text = _extract_json_text(data)
        renderer = "patroai_json_v1"
    elif intent.requested_format == "xlsx":
        data = _xlsx_bytes(normalized)
        semantic_text = _extract_xlsx_text(data)
        renderer = "patroai_xlsx_openpyxl_v1"
    elif intent.requested_format == "txt":
        data = normalized.encode("utf-8")
        semantic_text = data.decode("utf-8").strip()
        if not semantic_text:
            raise ArtifactValidationFailed("ARTIFACT_TXT_SEMANTIC_EMPTY")
        renderer = "patroai_text_v1"
    else:
        raise ArtifactFormatUnsupported("ARTIFACT_FORMAT_UNSUPPORTED")

    # Structured JSON is validated by parse/re-serialization; textual formats retain
    # the historical semantic probe to ensure the renderer did not drop source text.
    if intent.requested_format != "json":
        probe_source = normalized
        if intent.requested_format == "xlsx":
            semantic_rows: list[str] = []
            for line in normalized.splitlines():
                clean = line.strip()
                if not clean or _xlsx_separator_row(clean):
                    continue
                if "|" in clean:
                    cells = [cell.strip() for cell in clean.strip("|").split("|")]
                elif "\t" in clean:
                    cells = [cell.strip() for cell in clean.split("\t")]
                else:
                    cells = [clean]
                semantic_rows.append(" | ".join(cells))
            probe_source = "\n".join(semantic_rows)
        probe = re.sub(r"\s+", " ", probe_source).strip()[:80]
        if probe and probe not in re.sub(r"\s+", " ", semantic_text):
            raise ArtifactValidationFailed("ARTIFACT_SEMANTIC_MISMATCH")

    return ValidatedArtifactBytes(
        filename=filename,
        mime_type=intent.mime_type,
        data=data,
        sha256=hashlib.sha256(data).hexdigest(),
        semantic_text=semantic_text,
        renderer=renderer,
    )


def persist_validated_artifact(
    db: Session,
    *,
    settings: Settings,
    tenant_id: str,
    thread_id: str,
    created_by: str,
    validated: ValidatedArtifactBytes,
    source_message_sha256: str,
    source_response_message_id: str,
    agent_id: str,
    storage: BlobStorage | None = None,
) -> PersistedArtifact:
    storage = storage or build_blob_storage(settings)
    artifact_id = str(uuid.uuid4())
    key = f"{tenant_id}/{thread_id}/generated/{artifact_id}-{validated.filename}"
    provenance_key = key + ".provenance.json"
    metadata = {
        "artifact_id": artifact_id,
        "tenant_id": tenant_id,
        "thread_id": thread_id,
        "created_by": created_by,
        "agent_id": agent_id,
        "source_message_sha256": source_message_sha256,
        "source_response_message_id": source_response_message_id,
        "filename": validated.filename,
        "mime_type": validated.mime_type,
        "sha256": validated.sha256,
        "renderer": validated.renderer,
        "validated": True,
        "write_executed": True,
        "proposal_only": False,
    }

    artifact_created = False
    provenance_created = False
    try:
        final_bytes = validated.data
        if hashlib.sha256(final_bytes).hexdigest() != validated.sha256:
            raise ArtifactValidationFailed("ARTIFACT_FINAL_BYTES_HASH_MISMATCH")
        if validated.mime_type == DOCX_MIME:
            _extract_docx_text(final_bytes)
        elif validated.mime_type == PDF_MIME:
            _extract_pdf_text(final_bytes)
        elif validated.mime_type == PPTX_MIME:
            _extract_pptx_text(final_bytes)
        elif validated.mime_type == JSON_MIME:
            _extract_json_text(final_bytes)
        elif validated.mime_type == XLSX_MIME:
            _extract_xlsx_text(final_bytes)
        elif validated.mime_type in {TXT_MIME, MARKDOWN_MIME}:
            if not final_bytes.decode("utf-8").strip():
                raise ArtifactValidationFailed("ARTIFACT_FINAL_TEXT_EMPTY")
        artifact_created = storage.put_if_absent(
            key, final_bytes, content_type=validated.mime_type
        )
        stored_bytes = storage.get(key)
        if hashlib.sha256(stored_bytes).hexdigest() != validated.sha256:
            raise ArtifactValidationFailed("ARTIFACT_STORED_BYTES_HASH_MISMATCH")
        provenance_created = storage.put_if_absent(
            provenance_key,
            json.dumps(metadata, ensure_ascii=False, indent=2).encode("utf-8"),
            content_type="application/json",
        )

        row = Artifact(
            id=artifact_id,
            tenant_id=tenant_id,
            thread_id=thread_id,
            created_by=created_by,
            filename=validated.filename,
            mime_type=validated.mime_type,
            storage_key=key,
            sha256=validated.sha256,
            version=1,
        )
        db.add(row)
        db.commit()
    except BlobStorageError as exc:
        db.rollback()
        if artifact_created:
            try:
                storage.delete(key)
            except BlobStorageError:
                artifact_runtime_logger.warning(
                    "ARTIFACT_CLEANUP_FAILED %s",
                    json.dumps(
                        {"event": "artifact_cleanup_failed", "artifact_id": artifact_id, "path_name": key},
                        sort_keys=True,
                    ),
                )
        if provenance_created:
            try:
                storage.delete(provenance_key)
            except BlobStorageError:
                artifact_runtime_logger.warning(
                    "ARTIFACT_CLEANUP_FAILED %s",
                    json.dumps(
                        {"event": "artifact_cleanup_failed", "artifact_id": artifact_id, "path_name": provenance_key},
                        sort_keys=True,
                    ),
                )
        raise ArtifactStorageError("ARTIFACT_STORAGE_ERROR") from exc
    except Exception:
        db.rollback()
        if artifact_created:
            try:
                storage.delete(key)
            except BlobStorageError:
                artifact_runtime_logger.warning(
                    "ARTIFACT_CLEANUP_FAILED %s",
                    json.dumps(
                        {"event": "artifact_cleanup_failed", "artifact_id": artifact_id, "path_name": key},
                        sort_keys=True,
                    ),
                )
        if provenance_created:
            try:
                storage.delete(provenance_key)
            except BlobStorageError:
                artifact_runtime_logger.warning(
                    "ARTIFACT_CLEANUP_FAILED %s",
                    json.dumps(
                        {"event": "artifact_cleanup_failed", "artifact_id": artifact_id, "path_name": provenance_key},
                        sort_keys=True,
                    ),
                )
        raise

    result = PersistedArtifact(
        artifact=row,
        download_path=f"/api/v2/artifacts/{row.id}/download",
        provenance_path=provenance_key,
    )
    artifact_runtime_logger.info(
        "ARTIFACT_PERSISTED %s",
        json.dumps(
            {
                "event": "artifact_persisted",
                "artifact_id": row.id,
                "thread_id": thread_id,
                "agent_id": agent_id,
                "filename": row.filename,
                "mime_type": row.mime_type,
                "sha256": row.sha256,
                "renderer": validated.renderer,
                "download_path": result.download_path,
            },
            ensure_ascii=False,
            sort_keys=True,
        ),
    )
    return result


def artifact_payload(result: PersistedArtifact) -> dict[str, object]:
    row = result.artifact
    return {
        "id": row.id,
        "filename": row.filename,
        "mime_type": row.mime_type,
        "sha256": row.sha256,
        "version": row.version,
        "download_path": result.download_path,
        "created_at": row.created_at.isoformat() if row.created_at else None,
    }
