from __future__ import annotations

import hashlib
import io
import json
import logging
import re
import zipfile
from dataclasses import dataclass
from pathlib import Path

from sqlalchemy import select
from sqlalchemy.orm import Session

from ..config import Settings
from ..models import Attachment
from .document_ingestion_security import (
    ArchiveMemberTooLarge,
    ArchiveRequiredMemberInvalid,
    DocumentIngestionSecurityError,
    DocumentIngestionSecurityPolicy,
    InvalidXml,
    UnsafeXmlRejected,
    ArchiveSecurityLimits,
    default_document_ingestion_policy,
    inspect_zip_archive,
    parse_untrusted_xml,
    read_zip_member_bounded,
)


class DocumentContextError(RuntimeError):
    code = "DOCUMENT_CONTEXT_ERROR"


class DocumentStorageError(DocumentContextError):
    code = "DOCUMENT_STORAGE_ERROR"


class DocumentIntegrityError(DocumentContextError):
    code = "DOCUMENT_INTEGRITY_ERROR"


class DocumentExtractionUnsupported(DocumentContextError):
    code = "DOCUMENT_EXTRACTION_UNSUPPORTED"


class DocumentExtractionFailed(DocumentContextError):
    code = "DOCUMENT_EXTRACTION_FAILED"


_TRUNCATION_MARKER = "\n[document context truncated]"


@dataclass(frozen=True)
class ExtractedDocument:
    attachment_id: str
    filename: str
    mime_type: str
    text: str


@dataclass(frozen=True)
class DocumentSourceProvenance:
    attachment_id: str
    filename: str
    extraction_status: str
    source_chars: int
    provided_chars: int
    truncated: bool
    content_sha256: str


@dataclass(frozen=True)
class DocumentContextProvenance:
    available: bool
    sources: int
    source_ids: tuple[str, ...]
    extraction_status: str
    source_chars: int
    provided_chars: int
    per_source_truncated: bool
    aggregate_truncated: bool
    truncated: bool
    context_version: str = "1.1"
    source_provenance: tuple[DocumentSourceProvenance, ...] = ()


@dataclass(frozen=True)
class DocumentContextBundle:
    message: dict[str, str]
    provenance: DocumentContextProvenance
    errors: tuple[dict[str, str], ...]


_TEXT_MIME_TYPES = {
    "text/plain",
    "text/csv",
    "text/markdown",
    "application/json",
}
_PDF_MIME = "application/pdf"
_DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
_PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"

document_context_logger = logging.getLogger("uvicorn.error")


def _ingestion_policy_from_settings(settings: Settings) -> DocumentIngestionSecurityPolicy:
    return DocumentIngestionSecurityPolicy(
        archive=ArchiveSecurityLimits(
            max_entries=settings.document_ingestion_max_archive_entries,
            max_total_uncompressed_bytes=settings.document_ingestion_max_total_uncompressed_bytes,
            max_member_uncompressed_bytes=settings.document_ingestion_max_member_bytes,
        ),
        docx_max_xml_bytes=settings.document_ingestion_docx_max_xml_bytes,
    )


def _security_error_details(exc: DocumentContextError) -> dict[str, int | float | str]:
    details: dict[str, int | float | str] = {
        "ingestion_security_event": "document_ingestion_rejected",
        "error_code": str(exc) or exc.code,
    }
    cause = exc.__cause__
    if isinstance(cause, DocumentIngestionSecurityError):
        details.update(cause.safe_details())
    return details


def _safe_storage_path(settings: Settings, storage_key: str) -> Path:
    root = Path(settings.artifact_storage_path).resolve()
    target = (root / storage_key).resolve()
    if target != root and root not in target.parents:
        raise DocumentStorageError("DOCUMENT_STORAGE_PATH_INVALID")
    return target


def _normalise_text(value: str, *, max_chars: int | None = None) -> str:
    cleaned = value.replace("\x00", "")
    cleaned = re.sub(r"\r\n?", "\n", cleaned)
    cleaned = re.sub(r"[ \t]+", " ", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned).strip()
    if max_chars is not None and len(cleaned) > max_chars:
        cleaned = cleaned[:max_chars].rstrip() + _TRUNCATION_MARKER
    return cleaned


def _validate_magic(mime_type: str, raw: bytes) -> None:
    if mime_type == _PDF_MIME and not raw.startswith(b"%PDF-"):
        raise DocumentIntegrityError("DOCUMENT_MAGIC_MISMATCH")
    if mime_type == _DOCX_MIME:
        if not raw.startswith(b"PK\x03\x04"):
            raise DocumentIntegrityError("DOCUMENT_MAGIC_MISMATCH")
        try:
            with zipfile.ZipFile(io.BytesIO(raw)) as archive:
                names = set(archive.namelist())
        except zipfile.BadZipFile as exc:
            raise DocumentIntegrityError("DOCUMENT_CONTAINER_INVALID") from exc
        if "[Content_Types].xml" not in names or "word/document.xml" not in names:
            raise DocumentIntegrityError("DOCUMENT_CONTAINER_INVALID")
    if mime_type in _TEXT_MIME_TYPES and b"\x00" in raw[:4096]:
        raise DocumentIntegrityError("DOCUMENT_BINARY_CONTENT_REJECTED")


def _extract_text_plain(raw: bytes, *, mime_type: str) -> str:
    try:
        text = raw.decode("utf-8-sig")
    except UnicodeDecodeError as exc:
        raise DocumentExtractionFailed("DOCUMENT_TEXT_ENCODING_UNSUPPORTED") from exc
    if mime_type == "application/json":
        try:
            parsed = json.loads(text)
        except json.JSONDecodeError as exc:
            raise DocumentExtractionFailed("DOCUMENT_JSON_INVALID") from exc
        return json.dumps(parsed, ensure_ascii=False, indent=2)
    return text


def _extract_docx(
    raw: bytes,
    *,
    ingestion_policy: DocumentIngestionSecurityPolicy,
) -> str:
    try:
        with zipfile.ZipFile(io.BytesIO(raw)) as archive:
            inspect_zip_archive(
                archive,
                limits=ingestion_policy.archive,
                required_members=("word/document.xml",),
            )
            xml = read_zip_member_bounded(
                archive,
                "word/document.xml",
                max_bytes=ingestion_policy.docx_max_xml_bytes,
            )
    except ArchiveMemberTooLarge as exc:
        raise DocumentIntegrityError("DOCUMENT_DOCX_XML_TOO_LARGE") from exc
    except (DocumentIngestionSecurityError, zipfile.BadZipFile, KeyError) as exc:
        raise DocumentExtractionFailed("DOCUMENT_DOCX_INVALID") from exc

    try:
        root = parse_untrusted_xml(xml)
    except UnsafeXmlRejected as exc:
        raise DocumentExtractionFailed("DOCUMENT_DOCX_XML_UNSAFE") from exc
    except InvalidXml as exc:
        raise DocumentExtractionFailed("DOCUMENT_DOCX_XML_INVALID") from exc

    namespace = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    paragraphs: list[str] = []
    for paragraph in root.iter(namespace + "p"):
        parts = [node.text or "" for node in paragraph.iter(namespace + "t")]
        line = "".join(parts).strip()
        if line:
            paragraphs.append(line)
    return "\n".join(paragraphs)


def _extract_pptx(raw: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentExtractionUnsupported("DOCUMENT_PPTX_READER_UNAVAILABLE") from exc
    try:
        presentation = Presentation(io.BytesIO(raw))
        values: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = (shape.text or "").strip()
                    if text:
                        values.append(text)
        return "\n".join(values)
    except Exception as exc:
        raise DocumentExtractionFailed("DOCUMENT_PPTX_EXTRACTION_FAILED") from exc


def _extract_xlsx(raw: bytes) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError as exc:
        raise DocumentExtractionUnsupported("DOCUMENT_XLSX_READER_UNAVAILABLE") from exc
    try:
        workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        values: list[str] = []
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(value) for value in row if value is not None and str(value).strip()]
                if cells:
                    values.append(" | ".join(cells))
        workbook.close()
        return "\n".join(values)
    except Exception as exc:
        raise DocumentExtractionFailed("DOCUMENT_XLSX_EXTRACTION_FAILED") from exc


def _extract_pdf(raw: bytes, *, max_pages: int) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentExtractionUnsupported("DOCUMENT_PDF_READER_UNAVAILABLE") from exc
    try:
        reader = PdfReader(io.BytesIO(raw), strict=True)
        pages: list[str] = []
        for index, page in enumerate(reader.pages):
            if index >= max_pages:
                break
            pages.append(page.extract_text() or "")
        return "\n".join(pages)
    except Exception as exc:
        raise DocumentExtractionFailed("DOCUMENT_PDF_EXTRACTION_FAILED") from exc


def _extract_document_text_unbounded(
    *,
    filename: str,
    mime_type: str,
    raw: bytes,
    max_pdf_pages: int,
    ingestion_policy: DocumentIngestionSecurityPolicy | None = None,
) -> str:
    del filename  # filename is retained in the public signature for auditability/future adapters.
    _validate_magic(mime_type, raw)
    policy = ingestion_policy or default_document_ingestion_policy()
    if mime_type in _TEXT_MIME_TYPES:
        text = _extract_text_plain(raw, mime_type=mime_type)
    elif mime_type == _DOCX_MIME:
        text = _extract_docx(raw, ingestion_policy=policy)
    elif mime_type == _PPTX_MIME:
        text = _extract_pptx(raw)
    elif mime_type == _XLSX_MIME:
        text = _extract_xlsx(raw)
    elif mime_type == _PDF_MIME:
        text = _extract_pdf(raw, max_pages=max_pdf_pages)
    else:
        raise DocumentExtractionUnsupported("DOCUMENT_EXTRACTION_UNSUPPORTED")

    text = _normalise_text(text)
    if len(text) < 1:
        raise DocumentExtractionFailed("DOCUMENT_EXTRACTION_EMPTY")
    return text


def extract_document_text(
    *,
    filename: str,
    mime_type: str,
    raw: bytes,
    max_chars: int,
    max_pdf_pages: int,
    ingestion_policy: DocumentIngestionSecurityPolicy | None = None,
) -> str:
    """Backward-compatible extraction API.

    The legacy function still returns a prompt-ready string and therefore keeps the
    historical truncation marker. Provenance-aware callers use build_document_context(),
    where diagnostic marker characters are never counted as provided source characters.
    """
    text = _extract_document_text_unbounded(
        filename=filename,
        mime_type=mime_type,
        raw=raw,
        max_pdf_pages=max_pdf_pages,
        ingestion_policy=ingestion_policy,
    )
    if len(text) > max_chars:
        return text[:max_chars].rstrip() + _TRUNCATION_MARKER
    return text


def extract_document_blob(
    *,
    settings: Settings,
    filename: str,
    mime_type: str,
    raw: bytes,
) -> str:
    """Extract normalized source text using the same hardened ingestion policy.

    This helper is storage-agnostic and is used by the governed Knowledge Plane.
    It deliberately returns untruncated text so the caller can apply an aggregate
    budget across multiple knowledge sources.
    """
    return _extract_document_text_unbounded(
        filename=filename,
        mime_type=mime_type,
        raw=raw,
        max_pdf_pages=settings.document_context_max_pdf_pages,
        ingestion_policy=_ingestion_policy_from_settings(settings),
    )


def _status_for(*, successful_sources: int, errors: int, truncated: bool) -> str:
    if successful_sources == 0:
        return "failed" if errors else "none"
    if errors or truncated:
        return "partial"
    return "ready"


def build_document_context(
    db: Session,
    *,
    settings: Settings,
    tenant_id: str,
    thread_id: str,
) -> DocumentContextBundle | None:
    """Build document material plus truthful provenance for a canonical turn.

    source_chars is measured after extraction/normalisation but before character
    context limits. provided_chars counts only actual source characters supplied
    to the model; diagnostic marker text is excluded from that metric.
    """
    if not settings.document_context_enabled:
        return None

    rows = db.scalars(
        select(Attachment)
        .where(
            Attachment.tenant_id == tenant_id,
            Attachment.thread_id == thread_id,
        )
        .order_by(Attachment.created_at.asc(), Attachment.id.asc())
        .limit(settings.document_context_max_files)
    ).all()
    if not rows:
        return None

    remaining = max(0, int(settings.document_context_max_chars))
    blocks: list[str] = [
        "DOCUMENT CONTEXT — successfully extracted document content supplied below is "
        "available as source material for this turn. You may use the supplied content. "
        "Do not claim that no document content is available when content is supplied. "
        "Do not claim access to omitted or truncated portions, original file bytes, or "
        "attachments whose extraction failed."
    ]
    errors: list[dict[str, str]] = []
    source_provenance: list[DocumentSourceProvenance] = []
    source_ids: list[str] = []
    total_source_chars = 0
    total_provided_chars = 0
    any_per_file_truncated = False
    any_aggregate_truncated = False
    successful_sources = 0

    for attachment in rows:
        try:
            target = _safe_storage_path(settings, attachment.storage_key)
            raw = target.read_bytes()
            digest = hashlib.sha256(raw).hexdigest()
            if digest != attachment.sha256:
                raise DocumentIntegrityError("DOCUMENT_SHA256_MISMATCH")

            full_text = _extract_document_text_unbounded(
                filename=attachment.filename,
                mime_type=attachment.mime_type,
                raw=raw,
                max_pdf_pages=settings.document_context_max_pdf_pages,
                ingestion_policy=_ingestion_policy_from_settings(settings),
            )
            source_chars = len(full_text)
            per_limit = max(0, int(settings.document_context_max_chars_per_file))
            per_file_provided = full_text[:per_limit]
            per_file_truncated = source_chars > len(per_file_provided)

            provided_text = per_file_provided[:remaining]
            aggregate_cut_this_source = len(provided_text) < len(per_file_provided)
            provided_chars = len(provided_text)

            successful_sources += 1
            source_ids.append(attachment.id)
            total_source_chars += source_chars
            total_provided_chars += provided_chars
            any_per_file_truncated = any_per_file_truncated or per_file_truncated
            any_aggregate_truncated = any_aggregate_truncated or aggregate_cut_this_source

            content_sha = hashlib.sha256(full_text.encode("utf-8")).hexdigest()
            source_truncated = per_file_truncated or aggregate_cut_this_source
            source_provenance.append(
                DocumentSourceProvenance(
                    attachment_id=attachment.id,
                    filename=attachment.filename,
                    extraction_status="ready",
                    source_chars=source_chars,
                    provided_chars=provided_chars,
                    truncated=source_truncated,
                    content_sha256=content_sha,
                )
            )

            if provided_chars:
                marker = _TRUNCATION_MARKER if source_truncated else ""
                blocks.append(
                    f"\n--- attachment:{attachment.id} filename:{attachment.filename} "
                    f"mime:{attachment.mime_type} ---\n{provided_text}{marker}"
                )
                remaining -= provided_chars

            document_context_logger.info(
                "DOCUMENT_CONTEXT_EXTRACTED %s",
                json.dumps(
                    {
                        "event": "document_context_extracted",
                        "tenant_id": tenant_id,
                        "thread_id": thread_id,
                        "attachment_id": attachment.id,
                        "filename": attachment.filename,
                        "mime_type": attachment.mime_type,
                        "sha256": attachment.sha256,
                        "source_chars": source_chars,
                        "provided_chars": provided_chars,
                        "per_source_truncated": per_file_truncated,
                        "aggregate_truncated": aggregate_cut_this_source,
                        "truncated": source_truncated,
                    },
                    sort_keys=True,
                ),
            )
        except FileNotFoundError:
            code = "DOCUMENT_STORAGE_MISSING"
            errors.append({"attachment_id": attachment.id, "code": code})
            document_context_logger.warning(
                "DOCUMENT_CONTEXT_FAILED %s",
                json.dumps(
                    {
                        "event": "document_context_failed",
                        "tenant_id": tenant_id,
                        "thread_id": thread_id,
                        "attachment_id": attachment.id,
                        "filename": attachment.filename,
                        "mime_type": attachment.mime_type,
                        "sha256": attachment.sha256,
                        "error_code": code,
                    },
                    sort_keys=True,
                ),
            )
        except DocumentContextError as exc:
            code = str(exc) or exc.code
            errors.append({"attachment_id": attachment.id, "code": code})
            document_context_logger.warning(
                "DOCUMENT_CONTEXT_FAILED %s",
                json.dumps(
                    {
                        "event": "document_context_failed",
                        "tenant_id": tenant_id,
                        "thread_id": thread_id,
                        "attachment_id": attachment.id,
                        "filename": attachment.filename,
                        "mime_type": attachment.mime_type,
                        "sha256": attachment.sha256,
                        "error_code": code,
                    },
                    sort_keys=True,
                ),
            )

    # A source successfully extracted after the aggregate limit is exhausted is still
    # evidence of aggregate truncation, even though it contributes zero model chars.
    if successful_sources and any(item.provided_chars == 0 and item.source_chars > 0 for item in source_provenance):
        any_aggregate_truncated = True

    truncated = any_per_file_truncated or any_aggregate_truncated
    available = total_provided_chars > 0
    extraction_status = _status_for(
        successful_sources=successful_sources,
        errors=len(errors),
        truncated=truncated,
    )

    if errors:
        blocks.append("\nDOCUMENT EXTRACTION ERRORS:")
        for error in errors:
            blocks.append(f"- attachment:{error['attachment_id']} code:{error['code']}")

    provenance = DocumentContextProvenance(
        available=available,
        sources=successful_sources,
        source_ids=tuple(source_ids),
        extraction_status=extraction_status,
        source_chars=total_source_chars,
        provided_chars=total_provided_chars,
        per_source_truncated=any_per_file_truncated,
        aggregate_truncated=any_aggregate_truncated,
        truncated=truncated,
        source_provenance=tuple(source_provenance),
    )
    return DocumentContextBundle(
        message={"role": "system", "content": "\n".join(blocks)},
        provenance=provenance,
        errors=tuple(errors),
    )


def load_thread_documents(
    db: Session,
    *,
    settings: Settings,
    tenant_id: str,
    thread_id: str,
) -> tuple[list[ExtractedDocument], list[dict[str, str]]]:
    """Compatibility wrapper for existing callers/tests.

    New code should use build_document_context() for provenance.
    """
    bundle = build_document_context(
        db,
        settings=settings,
        tenant_id=tenant_id,
        thread_id=thread_id,
    )
    if bundle is None:
        return [], []

    docs: list[ExtractedDocument] = []
    for source in bundle.provenance.source_provenance:
        # Reconstruct prompt-visible text only. This wrapper is not a provenance API.
        prefix = f"--- attachment:{source.attachment_id} filename:{source.filename} "
        content = bundle.message["content"]
        at = content.find(prefix)
        if at < 0:
            continue
        start = content.find("\n", at)
        if start < 0:
            continue
        end = content.find("\n--- attachment:", start + 1)
        if end < 0:
            end = content.find("\nDOCUMENT EXTRACTION ERRORS:", start + 1)
        if end < 0:
            end = len(content)
        text = content[start + 1:end]
        docs.append(
            ExtractedDocument(
                attachment_id=source.attachment_id,
                filename=source.filename,
                mime_type="",
                text=text,
            )
        )
    return docs, list(bundle.errors)


def document_context_message(
    db: Session,
    *,
    settings: Settings,
    tenant_id: str,
    thread_id: str,
) -> dict | None:
    bundle = build_document_context(
        db,
        settings=settings,
        tenant_id=tenant_id,
        thread_id=thread_id,
    )
    return bundle.message if bundle is not None else None
